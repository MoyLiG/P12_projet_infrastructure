"""
Support de soutenance P12 (thème sombre charcoal).

Canvas charcoal sur TOUTES les slides, sections numérotées « 01 / », sommaire,
valeurs chiffrées en mono. Palette HYBRIDE — emerald (identité sport de P12,
accent principal) + violet (secondaire). Polices non installées → fallback
Segoe UI + Consolas (sûr sous Windows 11).

Produit le livrable OC `Le_Gall_Morgan_Option_B_1_support_062026.pptx`.

Usage : python docs/soutenance/build_pptx.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = Path(__file__).resolve().parent / "Le_Gall_Morgan_Option_B_1_support_062026.pptx"
TOTAL = 22

# ----- Palette (charcoal dark + hybride emerald/violet) ----------------
BG       = RGBColor(0x0F, 0x11, 0x17)   # charcoal — fond de toutes les slides
SURFACE  = RGBColor(0x16, 0x1D, 0x27)   # carte sur fond sombre
SURFACE2 = RGBColor(0x1E, 0x27, 0x33)   # carte élevée / encart
EMERALD  = RGBColor(0x10, 0xB9, 0x81)   # accent principal (sport)
EMERALD_L= RGBColor(0x34, 0xD3, 0x99)   # accent clair
VIOLET   = RGBColor(0x84, 0x05, 0xFF)   # secondaire
VIOLET_L = RGBColor(0xA8, 0x55, 0xF7)   # valeurs en mono
TXT      = RGBColor(0xF4, 0xF4, 0xF5)   # texte principal
MUTED    = RGBColor(0x94, 0xA3, 0xB8)   # texte muté (slate)
DIM      = RGBColor(0x5B, 0x6B, 0x7A)   # très muté
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

F_DISP = "Segoe UI"
F_BODY = "Segoe UI"
F_MONO = "Consolas"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


# ----- Helpers --------------------------------------------------------
def solid(shape, color, line=None, line_w=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w or 1)
    shape.shadow.inherit = False
    return shape


def rect(slide, x, y, w, h, color, line=None, line_w=None, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    return solid(shp, color, line, line_w)


def R(text, size, color, bold=False, font=F_BODY, italic=False):
    return (text, size, color, bold, font, italic)


def textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=6, line_spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if line_spacing:
            p.line_spacing = line_spacing
        for (text, size, color, bold, font, italic) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
            r.font.italic = italic
    return tb


def pill(slide, x, y, label, fill=EMERALD, txt=WHITE, w=None, font=F_BODY):
    w = w or (0.16 * len(label) + 0.5)
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.4))
    solid(p, fill)
    tf = p.text_frame
    tf.margin_top = 0; tf.margin_bottom = 0
    para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
    run = para.add_run(); run.text = label
    run.font.size = Pt(10.5); run.font.bold = True; run.font.color.rgb = txt
    run.font.name = font
    return w


prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    rect(s, -0.1, -0.1, 13.6, 7.7, BG)
    return s


def wordmark(s, x=0.6, y=0.4):
    """Petit logo texte récurrent (rappel du logo Sport Data Solution)."""
    rect(s, x, y + 0.02, 0.15, 0.15, EMERALD)
    textbox(s, x + 0.27, y - 0.04, 6, 0.3,
            [[R("SPORT DATA ", 10.5, WHITE, True), R("SOLUTION", 10.5, EMERALD_L, True)]])


def footer(s, n):
    textbox(s, 0.6, 7.06, 8, 0.3,
            [[R("POC Avantages Sportifs · Option B", 9, DIM, False)]])
    textbox(s, 11.7, 7.06, 1.1, 0.3,
            [[R(f"{n:02d} / {TOTAL}", 9, DIM, False, F_MONO)]], align=PP_ALIGN.RIGHT)


def content_title(s, kicker, title):
    """Bandeau de titre des slides de contenu (style « 0X / kicker »)."""
    wordmark(s)
    textbox(s, 0.6, 1.0, 11, 0.4, [[R(kicker.upper(), 12, EMERALD_L, True, F_MONO)]])
    textbox(s, 0.58, 1.35, 12.1, 0.9, [[R(title, 28, TXT, True, F_DISP)]])
    rect(s, 0.62, 2.12, 0.6, 0.05, EMERALD)


def card(s, x, y, w, h, bar=EMERALD, fill=SURFACE):
    rect(s, x, y, w, h, fill)
    rect(s, x, y, w, 0.07, bar)


# =====================================================================
# Slide 1 — Titre
# =====================================================================
s = new_slide()
rect(s, 0, 0, 0.26, 7.5, EMERALD)
# accents géométriques (dot/échelle façon dashboard)
rect(s, 10.9, 5.0, 2.0, 2.0, SURFACE)
rect(s, 11.2, 5.3, 1.4, 1.4, SURFACE2)
rect(s, 11.5, 5.6, 0.8, 0.8, EMERALD)
for i, hbar in enumerate([0.5, 0.9, 0.7, 1.2, 1.0]):
    rect(s, 10.95 + i * 0.42, 4.7 - hbar, 0.26, hbar,
         EMERALD if i % 2 else VIOLET_L)
pill(s, 0.9, 1.3, "OPENCLASSROOMS · DATA ENGINEER · OPTION B",
     fill=SURFACE2, txt=EMERALD_L, w=6.4, font=F_MONO)
textbox(s, 0.85, 2.3, 11.5, 1.6, [[R("POC Avantages Sportifs", 52, WHITE, True, F_DISP)]])
textbox(s, 0.9, 3.9, 11, 0.7,
        [[R("Pipeline ETL bout-en-bout — ", 22, MUTED, False),
          R("Sport Data Solution", 22, EMERALD_L, True)]])
textbox(s, 0.9, 5.5, 9, 1.2,
        [[R("Morgan Le Gall", 16, TXT, True)],
         [R("Juin 2026", 13, DIM, False, F_MONO)]], space_after=4)
# chips stack
textbox(s, 0.9, 6.55, 11, 0.3,
        [[R("Kestra   ", 11, EMERALD_L, True, F_MONO),
          R("PostgreSQL   ", 11, EMERALD_L, True, F_MONO),
          R("Python   ", 11, EMERALD_L, True, F_MONO),
          R("Great Expectations   ", 11, EMERALD_L, True, F_MONO),
          R("Docker", 11, EMERALD_L, True, F_MONO)]])

# =====================================================================
# Slide 2 — Sommaire
# =====================================================================
s = new_slide()
wordmark(s)
textbox(s, 0.58, 1.1, 11, 0.9, [[R("Sommaire", 34, WHITE, True, F_DISP)]])
rect(s, 0.62, 2.0, 0.6, 0.05, EMERALD)
sections = [
    ("01", "Contexte & mission", "deux avantages, trois objectifs"),
    ("02", "Architecture", "trois couches, un orchestrateur"),
    ("03", "Pipeline & qualité", "huit étapes, trois lignes de défense"),
    ("04", "Sécurité & données personnelles", "RGPD, pseudonymisation"),
    ("05", "Monitoring & résultats", "traçabilité, impact financier"),
    ("06", "Démo & industrialisation", "live + passage en production"),
]
for i, (num, t, sub) in enumerate(sections):
    col = 0.6 if i < 3 else 6.85
    y = 2.5 + (i % 3) * 1.4
    card(s, col, y, 5.9, 1.15, bar=(EMERALD if i % 2 == 0 else VIOLET))
    textbox(s, col + 0.3, y + 0.2, 1.2, 0.8,
            [[R(num, 30, EMERALD_L if i % 2 == 0 else VIOLET_L, True, F_MONO)]])
    textbox(s, col + 1.5, y + 0.2, 4.2, 0.5, [[R(t, 16, TXT, True, F_DISP)]])
    textbox(s, col + 1.5, y + 0.66, 4.2, 0.4, [[R(sub, 11, MUTED, False)]])
footer(s, 2)


# ---- Générateur de divider de section -------------------------------
def divider(num, title, subtitle, accent=EMERALD, accent_l=EMERALD_L):
    s = new_slide()
    rect(s, 0, 0, 0.26, 7.5, accent)
    wordmark(s)
    textbox(s, 0.85, 2.0, 8, 2.0, [[R(num, 150, accent, True, F_MONO)]])
    rect(s, 4.7, 2.55, 0.06, 2.4, SURFACE2)
    textbox(s, 5.0, 2.9, 7.8, 1.0, [[R("/", 40, accent_l, True, F_MONO),
                                     R("  " + title, 34, WHITE, True, F_DISP)]])
    textbox(s, 5.05, 4.0, 7.6, 0.8, [[R(subtitle, 15, MUTED, False)]], line_spacing=1.2)
    return s


# =====================================================================
# Section 01 — Contexte & mission
# =====================================================================
divider("01", "Contexte & mission", "Deux avantages sportifs à chiffrer avant d'industrialiser.")

# Slide 4 — Les deux avantages
s = new_slide()
content_title(s, "01 / Contexte", "La mission confiée par Juliette")
textbox(s, 0.6, 2.35, 12.0, 0.6,
        [[R("Encourager la pratique sportive des 161 salariés via deux avantages — "
            "et en mesurer le coût.", 14, MUTED, False)]], line_spacing=1.15)
def adv(x, tag, col, title, body):
    card(s, x, 3.15, 5.85, 3.0, bar=col)
    pill(s, x + 0.4, 3.45, tag, fill=col, w=1.9, font=F_MONO)
    textbox(s, x + 0.4, 4.15, 5.1, 0.7, [[R(title, 20, TXT, True, F_DISP)]])
    textbox(s, x + 0.4, 4.95, 5.1, 1.1, [[R(body, 13, MUTED, False)]], line_spacing=1.2)
adv(0.6, "PRIME", EMERALD, "+5 % du salaire brut",
    "Pour les salariés venant au bureau en mode actif (marche, course, vélo, "
    "trottinette). Vérifié par la distance domicile-travail.")
adv(6.85, "BIEN-ÊTRE", VIOLET, "5 jours offerts / an",
    "Pour les salariés réalisant au moins 15 activités physiques sur l'année, "
    "attestées par le flux de données sportives.")
footer(s, 4)

# Slide 5 — Objectifs
s = new_slide()
content_title(s, "01 / Objectifs", "Ce que le POC doit prouver")
objs = [
    ("01", "Faisabilité", "Un pipeline qui tourne de bout en bout, automatisé et reproductible."),
    ("02", "Données", "Identifier et structurer les données — dont des données RH sensibles."),
    ("03", "Impact financier", "Chiffrer le coût, et le recalculer si un paramètre change."),
]
for i, (num, t, b) in enumerate(objs):
    x = 0.6 + i * 4.05
    card(s, x, 2.6, 3.75, 3.3, bar=EMERALD if i != 1 else VIOLET)
    textbox(s, x + 0.35, 2.9, 2, 1, [[R(num, 38, EMERALD_L if i != 1 else VIOLET_L, True, F_MONO)]])
    textbox(s, x + 0.35, 4.0, 3.1, 0.7, [[R(t, 18, TXT, True, F_DISP)]])
    textbox(s, x + 0.35, 4.75, 3.1, 1.1, [[R(b, 12.5, MUTED, False)]], line_spacing=1.2)
footer(s, 5)

# =====================================================================
# Section 02 — Architecture
# =====================================================================
divider("02", "Architecture", "Un flux orchestré, trois couches de données.",
        accent=VIOLET, accent_l=VIOLET_L)

# Slide 7 — Architecture
s = new_slide()
content_title(s, "02 / Architecture", "Un flux orchestré, trois couches")
def abox(x, y, w, h, title, sub, bar=EMERALD):
    card(s, x, y, w, h, bar=bar)
    textbox(s, x + 0.18, y + 0.18, w - 0.36, 0.5, [[R(title, 13, TXT, True, F_DISP)]])
    if sub:
        textbox(s, x + 0.18, y + 0.66, w - 0.36, h - 0.7, [[R(sub, 10, MUTED, False)]], line_spacing=1.1)
def arrow(x, y, w=0.45):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.3))
    solid(a, EMERALD)
textbox(s, 0.6, 2.4, 3, 0.3, [[R("SOURCES", 10.5, VIOLET_L, True, F_MONO)]])
abox(0.6, 2.75, 3.0, 0.85, "Fichiers RH", "161 salariés · PII", bar=VIOLET)
abox(0.6, 3.75, 3.0, 0.85, "Pratiques sportives", "déclaratif", bar=VIOLET)
abox(0.6, 4.75, 3.0, 0.85, "Générateur Strava", "JSON façon API · 3948 act.", bar=VIOLET)
arrow(3.75, 3.9)
textbox(s, 4.4, 2.4, 4, 0.3, [[R("POSTGRESQL", 10.5, EMERALD_L, True, F_MONO)]])
abox(4.4, 2.75, 4.0, 2.85, "Base sécurisée",
     "raw (JSON Strava)  →  staging  →  marts\n\n+ audit (monitoring)\n+ cache (Google Maps)", bar=EMERALD)
arrow(8.55, 3.9)
textbox(s, 9.25, 2.4, 4, 0.3, [[R("VALIDATION & SORTIES", 10.5, EMERALD_L, True, F_MONO)]])
abox(9.25, 2.75, 3.45, 0.85, "Validation géo + qualité", "Google Maps · Great Expectations")
abox(9.25, 3.75, 3.45, 0.85, "Calcul des avantages", "prime · bien-être")
abox(9.25, 4.75, 1.62, 0.85, "Slack", "messages", bar=VIOLET)
abox(11.1, 4.75, 1.6, 0.85, "PowerBI", "KPIs", bar=VIOLET)
rect(s, 0.6, 5.95, 12.1, 0.65, SURFACE2)
textbox(s, 0.85, 6.08, 12, 0.45,
        [[R("KESTRA  ", 13, EMERALD_L, True, F_MONO),
          R("orchestre le tout — trigger mensuel, replay, monitoring, alertes", 12.5, TXT, False)]],
        anchor=MSO_ANCHOR.MIDDLE)
footer(s, 7)

# Slide 8 — Stack
s = new_slide()
content_title(s, "02 / Choix techniques", "La bonne brique pour chaque besoin")
rows = [
    ("Orchestration", "Kestra", "Replay de l'historique, monitoring, trigger cron"),
    ("Base de données", "PostgreSQL 16", "Rôles, RLS, audit — adapté aux PII RH"),
    ("Traitement", "Python + pandas", "Extraction, génération, calculs métier"),
    ("Qualité", "Great Expectations", "Tests déclaratifs, bloquants si échec"),
    ("Géocodage", "Google Maps API", "Distance réelle domicile-travail (+ cache)"),
    ("Restitution", "PowerBI", "KPIs, paramètre de taux dynamique"),
    ("Conteneurs", "Docker Compose", "Stack reproductible en une commande"),
]
y0 = 2.45
rect(s, 0.6, y0, 12.1, 0.45, SURFACE2)
for j, head in enumerate(["BRIQUE", "OUTIL", "POURQUOI"]):
    xs = [0.8, 3.5, 6.6][j]
    textbox(s, xs, y0 + 0.05, 6, 0.4, [[R(head, 10.5, EMERALD_L, True, F_MONO)]], anchor=MSO_ANCHOR.MIDDLE)
for i, (a, b, c) in enumerate(rows):
    y = y0 + 0.45 + i * 0.55
    if i % 2 == 0:
        rect(s, 0.6, y, 12.1, 0.55, SURFACE)
    textbox(s, 0.8, y + 0.07, 2.6, 0.45, [[R(a, 12, TXT, True)]], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, 3.5, y + 0.07, 3.0, 0.45, [[R(b, 12, EMERALD_L, True, F_MONO)]], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, 6.6, y + 0.07, 6.0, 0.45, [[R(c, 11.5, MUTED, False)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 8)

# =====================================================================
# Section 03 — Pipeline & qualité
# =====================================================================
divider("03", "Pipeline & qualité", "Huit étapes, et trois lignes de défense sur la donnée.")

# Slide 10 — Pipeline
s = new_slide()
content_title(s, "03 / Pipeline", "De la source à la restitution")
steps = [
    ("1", "Extract RH", "fichiers → base"),
    ("2", "Extract sport", "pratiques"),
    ("3", "Générer", "raw : JSON Strava"),
    ("4", "Transformer", "raw → staging"),
    ("5", "Valider géo", "Google Maps"),
    ("6", "Qualité", "Great Expectations"),
    ("7", "Calculer", "prime + bien-être"),
    ("8", "Notifier", "Slack"),
]
positions = [(0.6 + i * 3.05, 2.6) for i in range(4)] + [(0.6 + i * 3.05, 4.55) for i in range(4)]
for i, (num, t, b) in enumerate(steps):
    x, y = positions[i]
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.7), Inches(0.7))
    solid(circ, EMERALD if i % 2 == 0 else VIOLET)
    tf = circ.text_frame; tf.margin_top = 0; tf.margin_bottom = 0
    pr = tf.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
    rn = pr.add_run(); rn.text = num; rn.font.size = Pt(24); rn.font.bold = True
    rn.font.color.rgb = WHITE; rn.font.name = F_MONO
    textbox(s, x + 0.85, y + 0.0, 2.0, 0.5, [[R(t, 14, TXT, True, F_DISP)]])
    textbox(s, x + 0.85, y + 0.42, 2.0, 0.5, [[R(b, 10.5, MUTED, False)]])
rect(s, 0.6, 5.95, 12.1, 0.6, SURFACE2)
textbox(s, 0.85, 6.06, 12, 0.4,
        [[R("Échec d'une étape  →  ", 12, EMERALD_L, True, F_MONO),
          R("le pipeline s'arrête et alerte sur Slack. Jamais de calcul sur données douteuses.", 12, TXT, False)]],
        anchor=MSO_ANCHOR.MIDDLE)
footer(s, 10)

# Slide 11 — Qualité
s = new_slide()
content_title(s, "03 / Qualité", "Trois lignes de défense")
defs = [
    ("Contraintes SQL", "À l'entrée en base : distance ≥ 0, dates valides, transport "
     "dans une liste fermée. Rejet immédiat.", EMERALD),
    ("Great Expectations", "Règles métier déclaratives, bloquantes. A détecté un vrai "
     "bug : des activités datées dans le futur.", EMERALD_L),
    ("Validation géographique", "Déclaration « marche » mais domicile à 50 km ? Google "
     "Maps calcule la distance réelle → anomalie.", VIOLET),
]
for i, (t, b, col) in enumerate(defs):
    y = 2.45 + i * 1.3
    card(s, 0.6, y, 12.1, 1.15, bar=col, fill=SURFACE)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(y + 0.3), Inches(0.55), Inches(0.55))
    solid(circ, col)
    tf = circ.text_frame; tf.margin_top = 0; tf.margin_bottom = 0
    pr = tf.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
    rn = pr.add_run(); rn.text = str(i + 1); rn.font.size = Pt(20); rn.font.bold = True
    rn.font.color.rgb = BG; rn.font.name = F_MONO
    textbox(s, 1.8, y + 0.18, 10.6, 0.5, [[R(t, 17, TXT, True, F_DISP)]])
    textbox(s, 1.8, y + 0.64, 10.6, 0.45, [[R(b, 12.5, MUTED, False)]], line_spacing=1.1)
footer(s, 11)

# =====================================================================
# Section 04 — Sécurité & données personnelles
# =====================================================================
divider("04", "Sécurité & données personnelles", "Des données RH : cloisonnées, pseudonymisées, tracées.",
        accent=VIOLET, accent_l=VIOLET_L)

# Slide 13 — Sécurité & RGPD
s = new_slide()
content_title(s, "04 / Sécurité & RGPD", "Des données RH, donc cloisonnées")
points = [
    "3 rôles PostgreSQL : l'analyste ne voit jamais les noms ni salaires en clair",
    "Pseudonymisation : hash salé des identités, salaires en tranches",
    "Row Level Security + triggers d'audit (traçabilité RGPD)",
    "Aucun secret dans le code — variables d'environnement, .gitignore strict",
]
card(s, 0.6, 2.4, 7.6, 3.7, bar=EMERALD)
textbox(s, 0.95, 2.7, 7, 0.4, [[R("MESURES EN PLACE", 11.5, EMERALD_L, True, F_MONO)]])
textbox(s, 0.95, 3.25, 6.9, 2.7, [[R(p, 13.5, TXT, False)] for p in points],
        space_after=12, line_spacing=1.1)
card(s, 8.45, 2.4, 4.25, 3.7, bar=VIOLET, fill=SURFACE2)
textbox(s, 8.8, 2.7, 3.6, 0.4, [[R("AUDIT DE SÉCURITÉ", 11.5, VIOLET_L, True, F_MONO)]])
textbox(s, 8.8, 3.25, 3.6, 1.1, [[R("8 / 9", 50, EMERALD, True, F_MONO)]])
textbox(s, 8.8, 4.4, 3.6, 0.5, [[R("vulnérabilités corrigées", 13, TXT, False)]])
textbox(s, 8.8, 5.1, 3.6, 1.0,
        [[R("Bandit : 0 faille critique", 11.5, MUTED, False)],
         [R("Socket Docker retiré", 11.5, MUTED, False)],
         [R("Dépendances à jour", 11.5, MUTED, False)]], space_after=4)
footer(s, 13)

# Slide 14 — Cycle de vie PII
s = new_slide()
content_title(s, "04 / Données personnelles", "Le cycle de vie d'une donnée RH")
textbox(s, 0.6, 2.3, 12.0, 0.5,
        [[R("La même donnée, trois niveaux d'exposition selon la couche.", 14, MUTED, False)]])
def pii(x, tag, col, desc, example):
    w = 3.7
    card(s, x, 2.95, w, 2.55, bar=col)
    pill(s, x + 0.3, 3.2, tag, fill=col, w=1.5, font=F_MONO)
    textbox(s, x + 0.3, 3.85, w - 0.6, 1.0, [[R(desc, 12, MUTED, False)]], line_spacing=1.18)
    rect(s, x + 0.3, 4.95, w - 0.6, 0.45, SURFACE2)
    textbox(s, x + 0.42, 5.0, w - 0.8, 0.4, [[R(example, 10, VIOLET_L, False, F_MONO)]],
            anchor=MSO_ANCHOR.MIDDLE)
pii(0.6, "RAW", VIOLET, "Reflet brut du fichier RH. Nom, salaire, adresse en clair. "
    "Pipeline seul.", "Le Gall · 42 000 € · Lattes")
arrow_x = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.3), Inches(4.0), Inches(0.45), Inches(0.3))
solid(arrow_x, EMERALD)
pii(4.8, "STAGING", EMERALD, "Typé et hashé. PII en clair présentes mais cloisonnées "
    "par Row Level Security.", "Le Gall + hash a3f9c2…")
arrow_x2 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.5), Inches(4.0), Inches(0.45), Inches(0.3))
solid(arrow_x2, EMERALD)
pii(9.0, "MARTS", EMERALD_L, "Sortie analytique. Zéro PII : hash + tranche de salaire. "
    "Tout ce que voit PowerBI.", "a3f9c2… · BU Tech · 40-50 k€")
rect(s, 0.6, 5.8, 12.1, 0.75, SURFACE2)
textbox(s, 0.85, 5.92, 12, 0.5,
        [[R("etl_writer ", 12, EMERALD_L, True, F_MONO), R("écrit/lit les PII   ·   ", 12, TXT, False),
          R("analyst_reader / powerbi_reader ", 12, EMERALD_L, True, F_MONO),
          R("lisent marts — jamais de PII", 12, TXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 14)

# =====================================================================
# Section 05 — Monitoring & résultats
# =====================================================================
divider("05", "Monitoring & résultats", "Chaque run tracé ; l'impact financier chiffré.")

# Slide 16 — Monitoring
s = new_slide()
content_title(s, "05 / Monitoring", "Chaque exécution est tracée")
card(s, 0.6, 2.45, 5.7, 3.65, bar=EMERALD, fill=SURFACE2)
textbox(s, 0.95, 2.75, 5, 0.4, [[R("DERNIER RUN KESTRA", 11.5, EMERALD_L, True, F_MONO)]])
textbox(s, 0.95, 3.3, 5, 1.0, [[R("SUCCESS", 38, EMERALD, True, F_DISP)]])
textbox(s, 0.95, 4.35, 5, 0.5, [[R("12 / 12 tâches  ·  ~20 s", 14, TXT, False, F_MONO)]])
textbox(s, 0.95, 5.05, 5, 1.0,
        [[R("Volumétrie + durée par étape", 11.5, MUTED, False)],
         [R("Alerte Slack automatique si échec", 11.5, MUTED, False)],
         [R("Tests visibles + replay", 11.5, MUTED, False)]], space_after=5)
card(s, 6.55, 2.45, 6.15, 3.65, bar=VIOLET, fill=SURFACE)
textbox(s, 6.85, 2.7, 5.5, 0.4, [[R("audit.run_log (extrait)", 12, TXT, True, F_MONO)]])
logrows = [
    ("extract_rh", "161", "OK"),
    ("generate_activities", "3948", "OK"),
    ("transform_activities", "3948", "OK"),
    ("validate_quality", "3948", "OK"),
    ("compute_advantages", "68", "OK"),
    ("post_slack", "10", "OK"),
]
ly = 3.2
textbox(s, 6.85, ly, 3.2, 0.3, [[R("ÉTAPE", 9.5, MUTED, True, F_MONO)]])
textbox(s, 10.3, ly, 1.3, 0.3, [[R("LIGNES", 9.5, MUTED, True, F_MONO)]])
textbox(s, 11.8, ly, 0.9, 0.3, [[R("ST", 9.5, MUTED, True, F_MONO)]])
for i, (st, n, ok) in enumerate(logrows):
    y = ly + 0.42 + i * 0.48
    if i % 2 == 0:
        rect(s, 6.7, y - 0.04, 5.85, 0.46, SURFACE2)
    textbox(s, 6.85, y, 3.3, 0.4, [[R(st, 11, TXT, False, F_MONO)]], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, 10.3, y, 1.3, 0.4, [[R(n, 11, VIOLET_L, False, F_MONO)]], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, 11.8, y, 0.9, 0.4, [[R(ok, 11, EMERALD, True, F_MONO)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 16)

# Slide 17 — Résultats chiffrés
s = new_slide()
content_title(s, "05 / Résultats", "Ce que coûtent les avantages")
def kpi(x, y, w, number, label, col=EMERALD):
    card(s, x, y, w, 1.75, bar=col, fill=SURFACE)
    textbox(s, x + 0.3, y + 0.28, w - 0.5, 0.9, [[R(number, 34, col, True, F_MONO)]])
    textbox(s, x + 0.3, y + 1.15, w - 0.5, 0.5, [[R(label, 12, MUTED, False)]])
kpi(0.6, 2.5, 3.9, "172 482 €", "coût annuel de la prime", EMERALD)
kpi(4.7, 2.5, 3.9, "68 / 161", "éligibles à la prime", EMERALD_L)
kpi(8.8, 2.5, 3.9, "0", "déclaration suspecte", EMERALD_L)
kpi(0.6, 4.45, 3.9, "44", "éligibles bien-être", VIOLET)
kpi(4.7, 4.45, 3.9, "220", "jours bien-être accordés", VIOLET_L)
kpi(8.8, 4.45, 3.9, "3 948", "activités analysées", VIOLET_L)
textbox(s, 0.6, 6.35, 12, 0.5,
        [[R("Tous ces chiffres se recalculent en un clic si le taux de prime change.", 12.5, DIM, False, F_BODY, True)]])
footer(s, 17)

# =====================================================================
# Section 06 — Démo & industrialisation
# =====================================================================
divider("06", "Démo & industrialisation", "La démo live, puis le chemin vers la production.")

# Slide 19 — Démo
s = new_slide()
content_title(s, "06 / Démonstration", "Deux scénarios en direct")
def demo(x, tag, title, steps, col):
    card(s, x, 2.5, 5.85, 3.6, bar=col)
    pill(s, x + 0.4, 2.8, tag, fill=col, w=2.0, font=F_MONO)
    textbox(s, x + 0.4, 3.5, 5.1, 0.6, [[R(title, 18, TXT, True, F_DISP)]])
    textbox(s, x + 0.4, 4.25, 5.1, 1.7, [[R(stp, 13, MUTED, False)] for stp in steps],
            space_after=11, line_spacing=1.15)
demo(0.6, "SCÉNARIO A", "Changer le taux de prime",
     ["1. Passer le taux de 5 % à 7 % dans Kestra",
      "2. Relancer le pipeline",
      "3. Le coût total évolue dans PowerBI"], EMERALD)
demo(6.85, "SCÉNARIO B", "Insérer une activité live",
     ["1. Ajouter une nouvelle course",
      "2. Le message arrive dans Slack",
      "3. L'activité apparaît dans le reporting"], VIOLET)
footer(s, 19)

# Slide 20 — Scalabilité
s = new_slide()
content_title(s, "06 / Industrialisation", "Pensé pour passer en production")
left = [
    ("Reproductible", "Stack conteneurisée, démarrage en une commande"),
    ("Idempotent", "Rejouable sans doublons (Slack, calculs)"),
    ("Observable", "Tests visibles, alerte Slack, rapport GE exporté"),
]
right = [
    ("API Strava réelle", "Remplacer le générateur par les vraies données"),
    ("CI/CD", "Tests + audit sécurité à chaque commit"),
    ("Secret manager", "Rotation des secrets en production"),
]
textbox(s, 0.6, 2.45, 5, 0.4, [[R("DÉJÀ EN PLACE", 11.5, EMERALD_L, True, F_MONO)]])
textbox(s, 6.95, 2.45, 5, 0.4, [[R("PROCHAINES ÉTAPES", 11.5, VIOLET_L, True, F_MONO)]])
for i, (t, b) in enumerate(left):
    y = 2.95 + i * 1.1
    card(s, 0.6, y, 5.9, 0.95, bar=EMERALD)
    textbox(s, 0.95, y + 0.14, 5.3, 0.5, [[R(t, 15, TXT, True, F_DISP)]])
    textbox(s, 0.95, y + 0.56, 5.3, 0.4, [[R(b, 12, MUTED, False)]])
for i, (t, b) in enumerate(right):
    y = 2.95 + i * 1.1
    card(s, 6.8, y, 5.9, 0.95, bar=VIOLET)
    textbox(s, 7.15, y + 0.14, 5.3, 0.5, [[R(t, 15, TXT, True, F_DISP)]])
    textbox(s, 7.15, y + 0.56, 5.3, 0.4, [[R(b, 12, MUTED, False)]])
footer(s, 20)

# =====================================================================
# Slide 21 — Conclusion
# =====================================================================
s = new_slide()
rect(s, 0, 0, 0.26, 7.5, EMERALD)
wordmark(s)
pill(s, 0.85, 1.2, "CONCLUSION", fill=SURFACE2, txt=EMERALD_L, w=2.4, font=F_MONO)
textbox(s, 0.82, 1.9, 11.5, 1.2, [[R("Le POC valide les trois objectifs.", 32, WHITE, True, F_DISP)]])
checks = [
    "Faisable : pipeline complet, vert, en ~20 secondes",
    "Données maîtrisées : sensibles, cloisonnées, tracées",
    "Coût chiffré et pilotable : 172 482 € / an, recalculable à la demande",
]
textbox(s, 0.9, 3.5, 11, 2.0,
        [[R("✓  ", 16, EMERALD, True), R(c, 16, TXT, False)] for c in checks],
        space_after=14, line_spacing=1.1)
rect(s, 0.85, 5.95, 11.85, 0.7, EMERALD)
textbox(s, 1.1, 6.08, 11.4, 0.45,
        [[R("Recommandation : industrialiser, en commençant par l'intégration Strava.", 14, BG, True)]],
        anchor=MSO_ANCHOR.MIDDLE)
footer(s, 21)

# =====================================================================
# Slide 22 — Merci
# =====================================================================
s = new_slide()
rect(s, 10.6, 0, 2.73, 7.5, SURFACE)
rect(s, 11.0, 2.5, 1.9, 1.9, EMERALD)
rect(s, 11.25, 4.6, 0.26, 0.9, VIOLET_L)
rect(s, 11.65, 4.6, 0.26, 1.3, EMERALD)
rect(s, 12.05, 4.6, 0.26, 0.6, EMERALD_L)
wordmark(s)
textbox(s, 0.85, 2.6, 9.5, 1.2, [[R("Merci.", 50, WHITE, True, F_DISP)]])
textbox(s, 0.9, 3.85, 9, 0.7, [[R("Place à la discussion et à la démonstration.", 18, EMERALD_L, False)]])
textbox(s, 0.9, 5.1, 9, 0.5,
        [[R("Morgan Le Gall — POC Avantages Sportifs — Sport Data Solution", 13, DIM, False)]])
footer(s, 22)

n_slides = len(prs.slides._sldIdLst)
prs.save(str(OUT))
print(f"OK — {n_slides} slides écrites : {OUT}")
