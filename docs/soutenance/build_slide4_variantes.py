# -*- coding: utf-8 -*-
"""Genere schemas_p12_slide4_variantes.pptx : 2 variantes de la slide 4.
A = ER existant annote (faits / dimension). B = schema en etoile redessine.
"""
import copy
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SRC = "schemas_p12_corrige.pptx"
OUT = "schemas_p12_slide4_variantes.pptx"

# Palette extraite de la slide 4
BG        = RGBColor(0x0F,0x11,0x17)
SIDEBAR   = RGBColor(0x84,0x05,0xFF)
BODY      = RGBColor(0x16,0x1D,0x27)
EMERALD   = RGBColor(0x10,0xB9,0x81)   # dimension
MINT      = RGBColor(0x34,0xD3,0x99)
FACT      = RGBColor(0xF4,0x3F,0x5E)   # accent "fait" (rose, distinct d'audit/cache ambre)
VIEW      = RGBColor(0x38,0xBD,0xF8)   # vue (sky)
PK        = RGBColor(0xF5,0x9E,0x0B)
FK        = RGBColor(0x38,0xBD,0xF8)
FIELD     = RGBColor(0xF4,0xF4,0xF5)
MUTED     = RGBColor(0x94,0xA3,0xB8)
CAPTION   = RGBColor(0x5B,0x6B,0x7A)
LEGEND_BG = RGBColor(0x1E,0x27,0x33)
DARK_TXT  = RGBColor(0x0F,0x11,0x17)
WHITE     = RGBColor(0xFF,0xFF,0xFF)

def IN(v): return Emu(int(v*914400))

def no_line(sh): sh.line.fill.background()
def solid(sh, rgb):
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb; no_line(sh)

def txt(slide, x, y, w, h, runs, size=10, bold=False, align=PP_ALIGN.LEFT,
        font="Consolas", anchor=MSO_ANCHOR.MIDDLE):
    """runs = [(text, rgb, bold_override)] ou string."""
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left=Emu(0); tf.margin_right=Emu(0)
    tf.margin_top=Emu(0); tf.margin_bottom=Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    if isinstance(runs, str): runs=[(runs, FIELD, None)]
    for t,c,bo in runs:
        r=p.add_run(); r.text=t
        r.font.size=Pt(size); r.font.name=font
        r.font.bold = bold if bo is None else bo
        r.font.color.rgb=c
    return tb

def card(slide, x, y, w, h, title, header_rgb, rows, title_rgb=DARK_TXT, tag=None, tag_rgb=None):
    """Carte table : corps + bandeau header + lignes. rows=[(prefix,pcol,name,ncol)]."""
    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(x), IN(y), IN(w), IN(h))
    solid(body, BODY)
    hb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(x), IN(y), IN(w), IN(0.34))
    solid(hb, header_rgb)
    txt(slide, x+0.12, y+0.01, w-0.24, 0.32, [(title, title_rgb, True)], size=10, bold=True)
    ry = y+0.40
    for pre,pcol,name,ncol in rows:
        runs=[]
        if pre: runs.append((pre+" ", pcol, True))
        runs.append((name, ncol, False))
        txt(slide, x+0.12, ry, w-0.24, 0.24, runs, size=8)
        ry += 0.245
    if tag:
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(x), IN(y-0.30), IN(1.55), IN(0.26))
        solid(pill, tag_rgb)
        txt(slide, x, y-0.30, 1.55, 0.26, [(tag, DARK_TXT, True)], size=8.5, bold=True, align=PP_ALIGN.CENTER)
    return body

def arrow(slide, x1,y1,x2,y2, rgb=MUTED, w=1.5, dash=None):
    cn = slide.shapes.add_connector(2, IN(x1),IN(y1),IN(x2),IN(y2))  # 2 = straight
    cn.line.color.rgb=rgb; cn.line.width=Pt(w)
    if dash:
        ln=cn.line._get_or_add_ln(); d=ln.makeelement(qn('a:prstDash'),{'val':dash}); ln.append(d)
    return cn

def chrome(slide, sub_title, subtitle):
    """Fond + sidebar + bandeau titre commun."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(-0.1),IN(-0.1),IN(13.6),IN(7.7)); solid(bg,BG)
    sb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(0),IN(0),IN(0.26),IN(7.5)); solid(sb,SIDEBAR)
    sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(0.60),IN(0.50),IN(0.15),IN(0.15)); solid(sq,SIDEBAR)
    txt(slide,0.85,0.42,8.0,0.30,[("SPORT DATA SOLUTION",MUTED,True)],size=11,bold=True,font="Segoe UI")
    txt(slide,0.60,0.95,11.0,0.38,[(sub_title,WHITE,True)],size=20,bold=True,font="Segoe UI")
    txt(slide,0.58,1.30,12.2,0.40,[(subtitle,MUTED,False)],size=13,font="Segoe UI")
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(0.62),IN(1.95),IN(0.60),IN(0.05)); solid(accent,SIDEBAR)

def footer(slide, legend_runs, caption):
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(0.60),IN(6.70),IN(12.10),IN(0.44)); solid(bar,LEGEND_BG)
    txt(slide,0.85,6.74,12.0,0.36,legend_runs,size=9,font="Segoe UI")
    txt(slide,0.60,7.20,11.0,0.28,[(caption,CAPTION,False)],size=9,font="Segoe UI")
    txt(slide,11.40,7.20,1.40,0.28,[("POC Avantages Sportifs",CAPTION,False)],size=9,font="Segoe UI",align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------- duplication
def dup_slide(prs, src):
    layout = src.slide_layout
    new = prs.slides.add_slide(layout)
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    for el in list(src.shapes._spTree):
        if el.tag in (qn('p:nvGrpSpPr'), qn('p:grpSpPr')):
            continue
        new.shapes._spTree.append(copy.deepcopy(el))
    return new

prs = Presentation(SRC)
real4 = prs.slides[3]
varA = dup_slide(prs, real4)            # copie stylee de la vraie slide 4
varB = prs.slides.add_slide(real4.slide_layout)
for sh in list(varB.shapes):            # slide vierge pour l'etoile
    sh._element.getparent().remove(sh._element)

# ============================================================= VARIANTE A
# Annotation de l'ER existant : badges faits / dimension + callout.
def find(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle in sh.text:
            return sh
    return None

def retitle(slide, old, new):
    sh=find(slide, old)
    if sh:
        r=sh.text_frame.paragraphs[0].runs[0]; r.text=new
    return sh

def recolor_header(slide, title_needle, rgb):
    """Recolore le bandeau : parmi les formes au meme L/T, prend la plus basse (~0.34in)."""
    tb=find(slide, title_needle)
    if not tb: return
    L,T=tb.left,tb.top
    cand=[]
    for sh in slide.shapes:
        if sh.shape_type==1 and abs(sh.left-(L-int(0.12*914400)))<60000 and abs(sh.top-(T-int(0.02*914400)))<60000:
            cand.append(sh)
    if cand:
        bar=min(cand, key=lambda s:s.height)   # le bandeau est le plus fin
        solid(bar, rgb)

# titres + glyphes
retitle(varA, "marts.eligibility_prime", "marts.eligibility_prime  ★")
retitle(varA, "marts.eligibility_wellbeing", "marts.eligibility_wellbeing  ★")
retitle(varA, "staging.employees", "staging.employees  ◆")
recolor_header(varA, "marts.eligibility_prime", FACT)
recolor_header(varA, "marts.eligibility_wellbeing", FACT)
# staging.employees garde l'emeraude (dimension)

# callout "lecture dimensionnelle" dans la zone libre bas-droite
co = varA.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(7.42),IN(6.02),IN(5.28),IN(0.58))
co.fill.solid(); co.fill.fore_color.rgb=BODY; co.line.color.rgb=FACT; co.line.width=Pt(1)
txt(varA,7.58,6.05,5.0,0.52,
    [("★ ",FACT,True),("table de faits  ",FIELD,True),
     ("◆ ",EMERALD,True),("dimension conforme   ",FIELD,True),
     ("run_id + computed_at = dim. degenerees   ",MUTED,False),
     ("vue marts.kpi_by_bu = jointure etoile",VIEW,False)],
    size=8.5,font="Segoe UI",anchor=MSO_ANCHOR.MIDDLE)

# ============================================================= VARIANTE B
chrome(varB, "SCHÉMA 4 / MODÈLE DIMENSIONNEL",
        "Schéma en étoile : 2 tables de faits + 1 dimension conforme (grain : 1 salarié par run)")

dim_rows = [("PK",PK,"id_employee",FIELD),("",PK,"nom_hash",FIELD),
            ("",PK,"bu",FIELD),("",PK,"salaire_brut",FIELD),
            ("",PK,"distance_domicile_m",FIELD),("",PK,"moyen_deplacement",FIELD)]
prime_rows=[("PK·FK",FK,"id_employee",FIELD),("",PK,"is_eligible",FIELD),
            ("",PK,"prime_rate",FIELD),("",PK,"prime_amount",FIELD),
            ("",MUTED,"run_id  (deg.)",MUTED),("",MUTED,"computed_at  (deg.)",MUTED)]
well_rows =[("PK·FK",FK,"id_employee",FIELD),("",PK,"is_eligible",FIELD),
            ("",PK,"activity_count",FIELD),("",PK,"days_granted",FIELD),
            ("",MUTED,"run_id  (deg.)",MUTED),("",MUTED,"computed_at  (deg.)",MUTED)]
kpi_rows  =[("",VIEW,"bu",FIELD),("",PK,"nb_eligible_prime",FIELD),
            ("",PK,"total_prime_cost_eur",FIELD),("",PK,"total_wellbeing_days",FIELD)]

# Dimension au centre
card(varB, 5.15,2.62,3.00,2.05,"staging.employees",EMERALD,dim_rows,
     tag="◆ DIMENSION",tag_rgb=EMERALD)
# Faits gauche / droite
card(varB, 0.90,2.62,3.00,2.05,"marts.eligibility_prime",FACT,prime_rows,
     tag="★ FAIT",tag_rgb=FACT)
card(varB, 9.40,2.62,3.00,2.05,"marts.eligibility_wellbeing",FACT,well_rows,
     tag="★ FAIT",tag_rgb=FACT)
# Vue agregee bas-centre
card(varB, 5.15,5.20,3.00,1.18,"marts.kpi_by_bu",VIEW,kpi_rows,
     tag="vue · par BU",tag_rgb=VIEW)

# FK : faits -> dimension (lien id_employee)
arrow(varB, 3.90,3.55, 5.15,3.55, rgb=FK, w=2)
arrow(varB, 9.40,3.55, 8.15,3.55, rgb=FK, w=2)
txt(varB,4.00,3.18,1.05,0.26,[("FK id_employee",MUTED,False)],size=7,font="Consolas",align=PP_ALIGN.CENTER)
txt(varB,8.28,3.18,1.05,0.26,[("FK id_employee",MUTED,False)],size=7,font="Consolas",align=PP_ALIGN.CENTER)
# alimentation vers la vue (pointilles)
arrow(varB, 6.65,4.67, 6.65,5.20, rgb=VIEW, w=1.25, dash="dash")
arrow(varB, 2.40,4.67, 5.30,5.45, rgb=VIEW, w=1.0, dash="dash")
arrow(varB, 10.90,4.67, 8.00,5.45, rgb=VIEW, w=1.0, dash="dash")

footer(varB,
    [("★ ",FACT,True),("table de faits   ",FIELD,True),
     ("◆ ",EMERALD,True),("dimension conforme   ",FIELD,True),
     ("FK id_employee = rayon de l'étoile   ",MUTED,False),
     ("(deg.) run_id / computed_at = dimensions dégénérées",MUTED,False)],
    "Normalisé en amont (raw / staging) pour la qualité — dénormalisé en étoile en sortie (marts) pour la lecture PowerBI")

# ----------------------------------------------------- ne garder que A et B
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
for i in (0,1,2,3):       # supprime les 4 slides d'origine
    sldIdLst.remove(ids[i])

prs.save(OUT)
print("OK ->", OUT, "| slides:", len(prs.slides._sldIdLst))
