"""
Schémas de soutenance P12 — 4 diagrammes autonomes (thème tech dark charcoal).

Reprend EXACTEMENT la palette + les helpers de build_pptx.py pour rester
cohérent avec le support de soutenance. Produit un .pptx de 4 slides, chacune
un schéma plein cadre, à insérer dans la présentation existante :

  1. Architecture technique   — sources → PostgreSQL (3 couches) → sorties, orchestré par Kestra
  2. Flux ETL fonctionnel     — les 8 étapes en séquence + ce qui circule entre elles
  3. Cycle de vie / PII       — raw → staging → marts, exposition + rôles
  4. Modèle de données        — tables PostgreSQL par schéma + relations (FK)

Usage : python docs/soutenance/build_schemas.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

OUT = Path(__file__).resolve().parent / "schemas_p12.pptx"

# ----- Palette (identique à build_pptx.py) ----------------------------
BG       = RGBColor(0x0F, 0x11, 0x17)
SURFACE  = RGBColor(0x16, 0x1D, 0x27)
SURFACE2 = RGBColor(0x1E, 0x27, 0x33)
EMERALD  = RGBColor(0x10, 0xB9, 0x81)
EMERALD_L= RGBColor(0x34, 0xD3, 0x99)
VIOLET   = RGBColor(0x84, 0x05, 0xFF)
VIOLET_L = RGBColor(0xA8, 0x55, 0xF7)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
TXT      = RGBColor(0xF4, 0xF4, 0xF5)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
DIM      = RGBColor(0x5B, 0x6B, 0x7A)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

F_DISP = "Segoe UI"
F_BODY = "Segoe UI"
F_MONO = "Consolas"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


# ----- Helpers --------------------------------------------------------
def solid(shape, color, line=None, line_w=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w or 1)
    shape.shadow.inherit = False
    return shape


def rect(slide, x, y, w, h, color, line=None, line_w=None, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    return solid(shp, color, line, line_w)


def R(text, size, color, bold=False, font=F_BODY, italic=False):
    return (text, size, color, bold, font, italic)


def textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=6, line_spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if line_spacing:
            p.line_spacing = line_spacing
        for (text, size, color, bold, font, italic) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
            r.font.italic = italic
    return tb


def pill(slide, x, y, label, fill=EMERALD, txt=WHITE, w=None, font=F_BODY, h=0.4, size=10.5):
    w = w or (0.16 * len(label) + 0.5)
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    solid(p, fill)
    tf = p.text_frame
    tf.margin_top = 0; tf.margin_bottom = 0
    para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
    run = para.add_run(); run.text = label
    run.font.size = Pt(size); run.font.bold = True; run.font.color.rgb = txt
    run.font.name = font
    return w


def arrow(slide, x, y, w=0.5, color=EMERALD, h=0.3):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    return solid(a, color)


def down_arrow(slide, x, y, h=0.45, color=EMERALD, w=0.3):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    return solid(a, color)


def connector(slide, x1, y1, x2, y2, color=DIM, w=1.5, dash=False):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(w)
    cn.shadow.inherit = False
    if dash:
        # ligne pointillée via XML (python-pptx n'expose pas dash directement)
        from pptx.oxml.ns import qn
        ln = cn.line._get_or_add_ln()
        d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(d)
    return cn


prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    rect(s, -0.1, -0.1, 13.6, 7.7, BG)
    return s


def header(s, kicker, title, accent=EMERALD, accent_l=EMERALD_L):
    rect(s, 0, 0, 0.26, 7.5, accent)
    rect(s, 0.6, 0.5, 0.15, 0.15, accent)
    textbox(s, 0.85, 0.42, 8, 0.3, [[R("SPORT DATA ", 10.5, WHITE, True), R("SOLUTION", 10.5, accent_l, True)]])
    textbox(s, 0.6, 0.95, 11, 0.35, [[R(kicker.upper(), 12, accent_l, True, F_MONO)]])
    textbox(s, 0.58, 1.28, 12.2, 0.7, [[R(title, 26, TXT, True, F_DISP)]])
    rect(s, 0.62, 1.95, 0.6, 0.05, accent)


def footer(s, label):
    textbox(s, 0.6, 7.08, 9, 0.3, [[R(label, 9, DIM, False)]])
    textbox(s, 11.4, 7.08, 1.4, 0.3, [[R("POC Avantages Sportifs", 9, DIM, False)]], align=PP_ALIGN.RIGHT)


def node(s, x, y, w, h, title, sub=None, bar=EMERALD, fill=SURFACE, title_sz=12.5, sub_sz=9.5,
         title_col=TXT, mono_sub=False):
    rect(s, x, y, w, h, fill)
    rect(s, x, y, w, 0.06, bar)
    textbox(s, x + 0.16, y + 0.16, w - 0.32, 0.45, [[R(title, title_sz, title_col, True, F_DISP)]])
    if sub:
        textbox(s, x + 0.16, y + 0.16 + (title_sz / 24.0), w - 0.32, h - 0.5,
                [[R(sub, sub_sz, MUTED, False, F_MONO if mono_sub else F_BODY)]], line_spacing=1.05)
    return (x, y, w, h)


# =====================================================================
# SCHÉMA 1 — Architecture technique
# =====================================================================
s = new_slide()
header(s, "Schéma 1 / Architecture", "Architecture technique du pipeline")

# Zone SOURCES
textbox(s, 0.6, 2.25, 3, 0.3, [[R("SOURCES", 11, VIOLET_L, True, F_MONO)]])
node(s, 0.6, 2.6, 2.65, 0.78, "Fichiers RH (XLSX)", "161 salariés · PII", bar=VIOLET)
node(s, 0.6, 3.5, 2.65, 0.78, "Pratiques sportives", "déclaratif", bar=VIOLET)
node(s, 0.6, 4.4, 2.65, 0.78, "Générateur Strava", "JSON API-like · 3948", bar=VIOLET)

# Zone PostgreSQL (cœur)
textbox(s, 3.95, 2.25, 5, 0.3, [[R("POSTGRESQL 16  ·  base sécurisée", 11, EMERALD_L, True, F_MONO)]])
rect(s, 3.95, 2.6, 4.55, 3.55, SURFACE2)
rect(s, 3.95, 2.6, 4.55, 0.06, EMERALD)
# 3 couches
node(s, 4.15, 2.85, 4.15, 0.7, "raw", "payload JSON brut, jamais modifié", bar=VIOLET, fill=SURFACE, title_sz=12, mono_sub=False)
node(s, 4.15, 3.7, 4.15, 0.7, "staging", "typé · contraintes · hashé (RLS)", bar=EMERALD, fill=SURFACE, title_sz=12)
node(s, 4.15, 4.55, 4.15, 0.7, "marts", "sortie analytique, zéro PII", bar=EMERALD_L, fill=SURFACE, title_sz=12)
down_arrow(s, 6.1, 3.52, h=0.2, color=EMERALD)
down_arrow(s, 6.1, 4.37, h=0.2, color=EMERALD)
# annexes
node(s, 4.15, 5.4, 2.0, 0.62, "audit", "run_log", bar=AMBER, fill=SURFACE, title_sz=11, sub_sz=9)
node(s, 6.3, 5.4, 2.0, 0.62, "cache", "Google Maps", bar=AMBER, fill=SURFACE, title_sz=11, sub_sz=9)

arrow(s, 3.4, 3.75, w=0.45, color=EMERALD)

# Zone SORTIES / services
textbox(s, 9.2, 2.25, 4, 0.3, [[R("VALIDATION & SORTIES", 11, EMERALD_L, True, F_MONO)]])
node(s, 9.2, 2.6, 3.5, 0.78, "Google Maps API", "distance domicile-travail (+ cache)", bar=EMERALD, sub_sz=9)
node(s, 9.2, 3.5, 3.5, 0.78, "Great Expectations", "qualité déclarative, bloquante", bar=EMERALD, sub_sz=9)
node(s, 9.2, 4.4, 1.68, 0.78, "Slack", "alertes", bar=VIOLET, sub_sz=9)
node(s, 11.02, 4.4, 1.68, 0.78, "PowerBI", "KPIs", bar=VIOLET, sub_sz=9)
arrow(s, 8.6, 3.75, w=0.45, color=EMERALD)

# Bande Kestra
rect(s, 0.6, 6.35, 12.1, 0.6, EMERALD)
textbox(s, 0.85, 6.46, 12, 0.4,
        [[R("KESTRA  ", 13, BG, True, F_MONO),
          R("orchestre le tout — trigger mensuel · replay de l'historique · monitoring · alertes", 12, BG, True)]],
        anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Flux : sources → couches PostgreSQL → validation & restitution, le tout orchestré par Kestra")


# =====================================================================
# SCHÉMA 2 — Flux ETL fonctionnel (8 étapes)
# =====================================================================
s = new_slide()
header(s, "Schéma 2 / Flux ETL", "Le flux de traitement, étape par étape", accent=VIOLET, accent_l=VIOLET_L)

steps = [
    ("1", "Extract RH", "XLSX → raw.employees", "161 salariés", EMERALD),
    ("2", "Extract sport", "XLSX → raw.sports", "pratiques déclarées", EMERALD),
    ("3", "Générer", "→ raw.activities", "3948 JSON Strava", VIOLET),
    ("4", "Transformer", "raw → staging", "typage + hash PII", EMERALD),
    ("5", "Valider géo", "Google Maps", "distance réelle", EMERALD),
    ("6", "Qualité", "Great Expectations", "bloquant si échec", AMBER),
    ("7", "Calculer", "→ marts", "prime + bien-être", EMERALD_L),
    ("8", "Notifier", "Slack + audit", "trace du run", VIOLET),
]
# 2 rangées de 4, flux en serpentin
col_w, gap = 2.78, 0.32
x0 = 0.6
row_y = [2.5, 4.9]
positions = []
for r in range(2):
    for c in range(4):
        cc = c if r == 0 else 3 - c          # serpentin
        positions.append((x0 + cc * (col_w + gap), row_y[r]))

for i, (num, t, flow, note, col) in enumerate(steps):
    x, y = positions[i]
    rect(s, x, y, col_w, 1.55, SURFACE)
    rect(s, x, y, col_w, 0.06, col)
    # pastille numéro
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.22), Inches(0.5), Inches(0.5))
    solid(circ, col)
    tf = circ.text_frame; tf.margin_top = 0; tf.margin_bottom = 0
    pr = tf.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
    rn = pr.add_run(); rn.text = num; rn.font.size = Pt(18); rn.font.bold = True
    rn.font.color.rgb = WHITE; rn.font.name = F_MONO
    textbox(s, x + 0.82, y + 0.24, col_w - 0.9, 0.45, [[R(t, 14.5, TXT, True, F_DISP)]])
    rect(s, x + 0.2, y + 0.85, col_w - 0.4, 0.34, SURFACE2)
    textbox(s, x + 0.3, y + 0.88, col_w - 0.5, 0.3, [[R(flow, 10, EMERALD_L, False, F_MONO)]],
            anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, x + 0.2, y + 1.24, col_w - 0.4, 0.3, [[R(note, 10, MUTED, False)]])

# flèches horizontales rangée 1 (→)
for c in range(3):
    x = x0 + c * (col_w + gap) + col_w
    arrow(s, x + 0.01, row_y[0] + 0.6, w=gap - 0.02, color=VIOLET_L, h=0.3)
# flèche descente (col 4)
down_arrow(s, x0 + 3 * (col_w + gap) + col_w / 2 - 0.15, row_y[0] + 1.6, h=0.45, color=VIOLET_L)
# flèches horizontales rangée 2 (←)
for c in range(3):
    x = x0 + (3 - c) * (col_w + gap) - gap
    arrow(s, x + 0.01, row_y[1] + 0.6, w=gap - 0.02, color=VIOLET_L, h=0.3).rotation = 180

# bande règle d'arrêt
rect(s, 0.6, 6.65, 12.1, 0.5, SURFACE2)
textbox(s, 0.85, 6.72, 12, 0.4,
        [[R("Échec d'une étape  →  ", 11.5, AMBER, True, F_MONO),
          R("le pipeline s'arrête et alerte sur Slack. Jamais de calcul sur données douteuses.", 11.5, TXT, False)]],
        anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Pipeline orchestré par Kestra · run complet vert en ~20 s")


# =====================================================================
# SCHÉMA 3 — Cycle de vie d'une donnée / PII
# =====================================================================
s = new_slide()
header(s, "Schéma 3 / Données personnelles", "Cycle de vie d'une donnée RH")
textbox(s, 0.6, 2.2, 12, 0.45,
        [[R("La même donnée, trois niveaux d'exposition selon la couche traversée.", 13.5, MUTED, False)]])

def pii_card(x, tag, col, desc, example):
    w = 3.65
    rect(s, x, 2.85, w, 2.7, SURFACE)
    rect(s, x, 2.85, w, 0.06, col)
    pill(s, x + 0.3, 3.1, tag, fill=col, w=1.5, font=F_MONO)
    textbox(s, x + 0.3, 3.75, w - 0.6, 1.1, [[R(desc, 12, MUTED, False)]], line_spacing=1.2)
    rect(s, x + 0.3, 4.95, w - 0.6, 0.45, SURFACE2)
    textbox(s, x + 0.42, 5.0, w - 0.8, 0.4, [[R(example, 9.5, VIOLET_L, False, F_MONO)]],
            anchor=MSO_ANCHOR.MIDDLE)
    return w

pii_card(0.6, "RAW", VIOLET, "Reflet brut du fichier RH. Nom, salaire, adresse en clair. "
         "Accessible au pipeline seul.", "Le Gall · 42 000 € · Lattes")
arrow(s, 4.35, 4.0, w=0.5, color=EMERALD)
pii_card(4.85, "STAGING", EMERALD, "Typé et hashé. PII en clair présentes mais cloisonnées "
         "par Row Level Security.", "Le Gall + hash a3f9c2…")
arrow(s, 8.6, 4.0, w=0.5, color=EMERALD)
pii_card(9.1, "MARTS", EMERALD_L, "Sortie analytique. Zéro PII : hash + tranche de salaire. "
         "Tout ce que voit PowerBI.", "a3f9c2… · BU Tech · 40-50 k€")

# bande rôles
rect(s, 0.6, 5.9, 12.1, 1.0, SURFACE2)
textbox(s, 0.85, 6.02, 12, 0.3, [[R("QUI VOIT QUOI", 10.5, EMERALD_L, True, F_MONO)]])
textbox(s, 0.85, 6.36, 12, 0.45,
        [[R("etl_writer ", 12, EMERALD_L, True, F_MONO), R("écrit / lit les PII (raw + staging)", 12, TXT, False),
          R("       analyst_reader · powerbi_reader ", 12, VIOLET_L, True, F_MONO),
          R("lisent marts uniquement — jamais de PII", 12, TXT, False)]],
        anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Pseudonymisation (hash salé) + Row Level Security + triggers d'audit = conformité RGPD")


# =====================================================================
# SCHÉMA 4 — Modèle de données (ER simplifié)
# =====================================================================
s = new_slide()
header(s, "Schéma 4 / Modèle de données", "Tables PostgreSQL et leurs relations",
       accent=VIOLET, accent_l=VIOLET_L)

def table_box(x, y, w, name, cols, bar=EMERALD):
    h = 0.42 + 0.245 * len(cols)
    rect(s, x, y, w, h, SURFACE)
    rect(s, x, y, w, 0.34, bar)
    textbox(s, x + 0.12, y + 0.02, w - 0.2, 0.32, [[R(name, 10.5, BG, True, F_MONO)]],
            anchor=MSO_ANCHOR.MIDDLE)
    for i, (c, kind) in enumerate(cols):
        cy = y + 0.36 + i * 0.245
        if kind == "pkfk":
            ccol, prefix = AMBER, "PK·FK "
        elif kind == "pk":
            ccol, prefix = AMBER, "PK "
        elif kind == "fk":
            ccol, prefix = VIOLET_L, "FK "
        else:
            ccol, prefix = MUTED, "   "
        textbox(s, x + 0.12, cy, w - 0.2, 0.24,
                [[R(prefix, 7.5, ccol, True, F_MONO), R(c, 8.5, TXT if kind != "" else MUTED, False, F_MONO)]],
                anchor=MSO_ANCHOR.MIDDLE)
    return (x, y, w, h)

# en-têtes de couches
for lx, lab, lc in [(0.6, "raw", VIOLET_L), (4.55, "staging", EMERALD_L),
                    (8.5, "marts", EMERALD_L), (11.4, "audit / cache", AMBER)]:
    textbox(s, lx, 2.2, 1.85, 0.3, [[R(lab.upper(), 10.5, lc, True, F_MONO)]])

# raw
t_remp = table_box(0.6, 2.55, 2.55, "raw.employees_xlsx",
                   [("id_salarie", ""), ("nom, prenom", ""), ("salaire_brut", ""),
                    ("adresse_domicile", ""), ("moyen_deplacement", "")], bar=VIOLET)
t_ract = table_box(0.6, 4.7, 2.55, "raw.activities",
                   [("activity_id", "pk"), ("athlete_id", ""), ("payload JSONB", "")], bar=VIOLET)

# staging
t_semp = table_box(4.55, 2.55, 2.75, "staging.employees",
                   [("id_employee", "pk"), ("nom_hash", ""), ("salaire_brut", ""),
                    ("distance_domicile_m", ""), ("is_declaration_suspect", "")], bar=EMERALD)
t_sact = table_box(4.55, 4.7, 2.75, "staging.activities",
                   [("id_activity", "pk"), ("source_activity_id", "fk"),
                    ("id_employee", "fk"), ("sport_type, distance_m", "")], bar=EMERALD)

# marts
t_prime = table_box(8.5, 2.55, 2.7, "marts.eligibility_prime",
                    [("id_employee", "pkfk"), ("is_eligible", ""),
                     ("prime_amount", ""), ("run_id", "")], bar=EMERALD_L)
t_well = table_box(8.5, 4.55, 2.7, "marts.eligibility_wellbeing",
                   [("id_employee", "pkfk"), ("activity_count", ""),
                    ("days_granted", ""), ("run_id", "")], bar=EMERALD_L)

# audit / cache
t_log = table_box(11.4, 2.55, 1.85, "audit.run_log",
                  [("run_id", ""), ("step_name", ""), ("status", "")], bar=AMBER)
t_cache = table_box(11.4, 4.2, 1.85, "cache.gmaps",
                    [("address_hash", "pk"), ("mode", "pk"), ("distance_m", "")], bar=AMBER)

# relations (FK)
def mid_r(t): return (t[0] + t[2], t[1] + t[3] / 2)
def mid_l(t): return (t[0], t[1] + t[3] / 2)

# raw → staging (alimentation, pointillé)
x1, y1 = mid_r(t_remp); x2, y2 = mid_l(t_semp)
connector(s, x1, y1, x2, y2, color=DIM, w=1.5, dash=True)
x1, y1 = mid_r(t_ract); x2, y2 = mid_l(t_sact)
connector(s, x1, y1, x2, y2, color=DIM, w=1.5, dash=True)
# staging.activities.source_activity_id → raw.activities (FK)
connector(s, t_sact[0], t_sact[1] + 0.5, t_ract[0] + t_ract[2], t_ract[1] + 0.5, color=VIOLET_L, w=1.75)
# staging.employees → staging.activities (FK id_employee)
connector(s, t_semp[0] + 0.4, t_semp[1] + t_semp[3], t_sact[0] + 0.4, t_sact[1], color=EMERALD, w=1.75)
# staging.employees → marts.prime / wellbeing (FK)
x1, y1 = mid_r(t_semp)
connector(s, x1, y1, t_prime[0], t_prime[1] + 0.5, color=EMERALD, w=1.75)
connector(s, x1, y1, t_well[0], t_well[1] + 0.5, color=EMERALD, w=1.75)

# légende
rect(s, 0.6, 6.75, 12.1, 0.42, SURFACE2)
textbox(s, 0.85, 6.79, 12, 0.34,
        [[R("PK", 9, AMBER, True, F_MONO), R(" clé primaire    ", 9.5, MUTED, False),
          R("FK", 9, VIOLET_L, True, F_MONO), R(" clé étrangère    ", 9.5, MUTED, False),
          R("———", 9, EMERALD, True, F_MONO), R(" référence    ", 9.5, MUTED, False),
          R("- - -", 9, DIM, True, F_MONO), R(" alimentation (raw → staging)", 9.5, MUTED, False)]],
        anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Schémas raw / staging / marts + audit & cache — lignage tracé de la source à la restitution")


n_slides = len(prs.slides._sldIdLst)
prs.save(str(OUT))
print(f"OK — {n_slides} slides écrites : {OUT}")
